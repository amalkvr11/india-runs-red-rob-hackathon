from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation(r'D:\red_rob_hackathon\Idea Submission Template _ Redrob_Filled.pptx')

slide = prs.slides[6]  # Slide 7 (0-indexed)

# Add a text box for architecture content
from pptx.util import Emu
left = Emu(600000)
top = Emu(1400000)
width = Emu(8000000)
height = Emu(3600000)

txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.word_wrap = True

# Title
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Data Flow"
run.font.size = Pt(16)
run.font.bold = True
run.font.color.rgb = RGBColor(0x00, 0x96, 0x88)

# Data flow
lines = [
    "",
    "candidates.jsonl (100K profiles, 487 MB)",
    "         │",
    "         ▼",
    "┌────────────────────────────────────────────────────────┐",
    "│  Pipeline Orchestrator — 8 Scorers (weighted 0.05–0.22) │",
    "│  Title/Role │ Skills │ Career Quality │ Experience     │",
    "│  Statement  │ Behavioral │ Location │ Education        │",
    "└────────────────────────────────┬───────────────────────┘",
    "                                 │",
    "                                 ▼",
    "┌────────────────────────────────────────────────────────┐",
    "│  Honeypot Detector — 12+ anomaly signals                │",
    "│  (expert skill threshold, YOE mismatch, future dates,   │",
    "│   assessment contradictions, excessive skills/careers)  │",
    "│  → penalty up to 0.6× on final score                    │",
    "└────────────────────────────────┬───────────────────────┘",
    "                                 │",
    "                                 ▼",
    "┌────────────────────────────────────────────────────────┐",
    "│  Final Score = Weighted_Sum × (1 − Honeypot_Penalty)   │",
    "└────────────────────────────────┬───────────────────────┘",
    "                                 │",
    "                                 ▼",
    "┌────────────────────────────────────────────────────────┐",
    "│  Sort (score desc) → Top 100 → submission.csv          │",
    "└────────────────────────────────────────────────────────┘",
]

for i, line in enumerate(lines):
    if i == 0:
        continue  # first para already used for title
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = line
    if '┌' in line or '│' in line or '└' in line or '─' in line or '┬' in line or '┤' in line or '├' in line or '┴' in line or '▼' in line or '→' in line:
        run.font.name = 'Consolas'
        run.font.size = Pt(8)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    else:
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# Add a sub-section for interfaces
p = tf.add_paragraph()
run = p.add_run()
run.text = ""
p = tf.add_paragraph()
run = p.add_run()
run.text = "Additional Interfaces"
run.font.size = Pt(14)
run.font.bold = True
run.font.color.rgb = RGBColor(0x00, 0x96, 0x88)

interfaces = [
    "• Streamlit Dashboard (app.py) — interactive exploration & visualization",
    "• FastAPI Backend (api/server.py) — REST API for ranking & results",
    "• Vue 3 Frontend (frontend/) — full SPA with PrimeVU components",
    "• Git LFS — version control for large dataset (487 MB)",
]
for line in interfaces:
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = line
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

output = r'D:\red_rob_hackathon\Idea Submission Template _ Redrob_v2.pptx'
prs.save(output)
print(f'Slide 7 updated and saved to: {output}')
