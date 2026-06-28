from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re

prs = Presentation(r'D:\red_rob_hackathon\Idea Submission Template _ Redrob.pptx')

def find_placeholder(slide, marker):
    """Find a shape containing placeholder text marker and return it."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                full = para.text.strip()
                if marker in full:
                    return shape
    return None

def set_shape_text(shape, text, size=12, bold=False, color=RGBColor(0x33, 0x33, 0x33)):
    """Replace all text in a shape."""
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

def append_to_shape(shape, text, size=11, bold=False, color=RGBColor(0x33, 0x33, 0x33), new_para=True):
    """Append text to an existing paragraph or add a new paragraph."""
    tf = shape.text_frame
    if new_para:
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[-1]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

# ── Slide 1: Title ──
for shape in prs.slides[0].shapes:
    if shape.has_text_frame:
        full = shape.text_frame.text.strip()
        if full == 'Team Name :':
            set_shape_text(shape, 'Team Name : Redrob Ranker')
        elif full == 'Problem Statement :':
            set_shape_text(shape, 'Problem Statement : Rank 100K+ candidates for a Senior AI Engineer role at a Series A startup using multi-dimensional scoring with honeypot fraud detection.')
        elif full == 'Team Leader Name :':
            set_shape_text(shape, 'Team Leader Name : Amal K V')

# ── Slide 2: Solution Overview ──
shape = find_placeholder(prs.slides[1], 'What is your proposed solution')
if shape:
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "A multi-dimensional scoring engine that evaluates 100K+ candidate profiles "
        "against a job description across 8 weighted dimensions (Title/Role, Skills, "
        "Career Quality, Experience, Statement, Behavioral, Location, Education), then "
        "applies a honeypot fraud-detection penalty. The system outputs the top 100 "
        "ranked candidates with explainable reasoning."
    )
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

shape2 = find_placeholder(prs.slides[1], 'What differentiates')
if shape2:
    tf = shape2.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "Unlike traditional keyword-matching ATS, our approach uses:\n"
        "• 8-dimensional weighted scoring calibrated to the JD\n"
        "• Title tier system with ML-description rescue for misclassified roles\n"
        "• Honeypot detection (12+ anomaly signals) to disqualify synthetic profiles\n"
        "• Explainable output per candidate — each dimension has a human-readable reasoning string"
    )
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# ── Slide 3: JD Understanding ──
shape = find_placeholder(prs.slides[2], 'key requirements extracted')
if shape:
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "Key JD requirements extracted:\n"
        "• AI/ML engineering title (tier 1.0)\n"
        "• Skills: embeddings, vector search, NLP, LLMs, ranking metrics\n"
        "• 3-10 years experience with production ML exposure\n"
        "• Startup mindset, fast-paced, building from scratch\n"
        "• India tech hub location preferred\n"
        "• Consulting firm background is a disqualifier"
    )
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

shape2 = find_placeholder(prs.slides[2], 'most important')
if shape2:
    tf = shape2.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "Beyond keyword matching, we evaluate:\n"
        "• Title/Role (0.22): Current title tier + ML keyword density in descriptions\n"
        "• Skills (0.18): Skill group relevance × proficiency × endorsements × duration\n"
        "• Career Quality (0.18): Tenure stability, seniority progression, production exposure\n"
        "• Statement (0.10): ML intent + startup affinity + impact language in summary\n"
        "• Behavioral (0.09): Recency, response rate, verification, salary alignment"
    )
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# ── Slide 4: Ranking Methodology ──
shape = find_placeholder(prs.slides[3], 'How does your system retrieve')
if shape:
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "Retrieval: Full scan of 100K JSONL candidates (single-pass, streaming parser).\n\n"
        "Scoring pipeline (per candidate):\n"
        "1. Eight independent scorers compute sub-scores (title_role, skills, career_quality, etc.)\n"
        "2. Weighted sum: final = sum(sub_score[i] × weight[i])\n"
        "3. Honeypot detection: 12+ anomaly signals flag synthetic profiles\n"
        "4. Final score = weighted_sum × (1 − honeypot_penalty)\n"
        "5. Sort by score desc → select top 100"
    )
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

shape2 = find_placeholder(prs.slides[3], 'What models, algorithms')
if shape2:
    tf = shape2.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "Algorithms & heuristics:\n"
        "• Title tier system: 80+ normalized titles mapped to scores (0.0–1.0)\n"
        "• Skills taxonomy: 7 skill groups with proficiency × endorsements × duration formula\n"
        "• Career quality: tenure stability, seniority progression detection, production keyword density\n"
        "• Honeypot rules: expert skill threshold, YOE mismatch, future dates, assessment contradictions\n"
        "• All pure Python — no ML models, no external API calls, CPU-only"
    )
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# ── Slide 5: Explainability & Data Validation ──
shape = find_placeholder(prs.slides[4], 'How are ranking decisions explained')
if shape:
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "Each of the top 100 candidates includes a human-readable reasoning string:\n"
        'e.g., "Strong title: Senior AI Engineer | Skills: Python, PyTorch, NLP | '
        '@Google | Edu: M.Tech AI from IIT | Statement shows strong ML intent | '
        'Score=0.8921 | (title_role=0.950; skills=0.870; career_quality=0.760; ...)"\n\n'
        "All 8 dimension sub-scores are exposed for debugging and transparency."
    )
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

shape2 = find_placeholder(prs.slides[4], 'do you prevent hallucinations')
if shape2:
    tf = shape2.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "No LLM/AI generation — all reasoning is template-based from actual profile data. "
        "Each reasoning string is constructed from explicit field values (title, company, "
        "education, skills, computed sub-scores), eliminating hallucinations entirely."
    )
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

shape3 = find_placeholder(prs.slides[4], 'inconsistent, low-quality, or suspicious')
if shape3:
    tf = shape3.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "Honeypot detection system with 12+ signals:\n"
        "• Many expert skills with low endorsements → impossible profile\n"
        "• YOE > career history by 3+ years → fabrication\n"
        "• Future last_active_date → data error\n"
        "• Assessment contradiction (expert skill + 0 score)\n"
        "• Too many skills (>30) / too many career entries (>15) → script-generated\n"
        "• Missing required fields\n"
        "Penalty up to 0.6× on final score — honeypots are demoted below legitimate profiles."
    )
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# ── Slide 6: End-to-End Workflow ──
shape = find_placeholder(prs.slides[5], 'What is the complete workflow')
if shape:
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "1. INPUT: candidates.jsonl (100K profiles, ~465 MB)\n"
        "2. LOAD: Streaming JSONL parser (1 pass, ~4 GB peak memory)\n"
        "3. SCORE (per candidate):\n"
        "   a. 8 scorers fire in sequence (title_role → skills → career_quality → experience → statement → behavioral → location → education)\n"
        "   b. Honeypot detector runs (12+ signals, up to 0.6× penalty)\n"
        "   c. weighted_sum × (1 − honeypot_penalty) = final score\n"
        "4. SORT: All 100K scored candidates sorted by score descending\n"
        "5. OUTPUT: Top 100 → submission.csv (candidate_id, rank, score, reasoning)\n"
        "6. VALIDATE: submission_metadata.yaml describes dimensions + scoring approach"
    )
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# ── Slide 7: System Architecture ──
# Add architecture text below the slide title
for shape in prs.slides[6].shapes:
    if shape.has_text_frame and shape.text_frame.text.strip() == '':
        tf = shape.text_frame
        # Check if it's a content box (not the title)
        if shape.top > Inches(1):
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = (
                "INPUT: candidates.jsonl (100K profiles, ~465 MB)\n"
                "        │\n"
                "        ▼\n"
                "Pipeline Orchestrator — 8 Scorers:\n"
                "  Title/Role │ Skills │ Career Quality │ Experience\n"
                "  Statement  │ Behavioral │ Location │ Education\n"
                "        │\n"
                "        ▼\n"
                "Honeypot Detector — 12+ anomaly signals (penalty up to 0.6x)\n"
                "        │\n"
                "        ▼\n"
                "Score = Weighted_Sum × (1 − Honeypot_Penalty)\n"
                "        │\n"
                "        ▼\n"
                "Sort → Top 100 → submission.csv\n"
                "\n"
                "Additional interfaces:\n"
                "- Streamlit Dashboard (app.py) — interactive exploration\n"
                "- FastAPI Backend (api/server.py) — REST API\n"
                "- Vue.js Frontend (frontend/) — full web UI"
            )
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            run.font.name = 'Consolas'
            break

# ── Slide 8: Results & Performance ──
shape = find_placeholder(prs.slides[7], 'What results or insights')
if shape:
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "Results from 100K candidate profiles:\n"
        "• Top-ranked candidates show strong AI/ML title match + skills alignment + quality career history\n"
        "• Honeypot profiles consistently scored near the bottom (penalized 0.4–0.6×)\n"
        "• Top 100 dominated by AI Engineers, ML Engineers, Data Scientists with relevant skill sets\n"
        "• Location: majority in Bangalore, Hyderabad, Pune, Gurgaon (preferred tech hubs)\n"
        "• Experience sweet spot: 4–10 years within ideal 3–10 year range"
    )
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

shape2 = find_placeholder(prs.slides[7], 'runtime and compute constraints')
if shape2:
    tf = shape2.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "Performance metrics:\n"
        "• ~64 seconds for full 100K pipeline (single-threaded, CPU only)\n"
        "• ~4 GB peak RAM (well within 16 GB limit)\n"
        "• No GPU, no network calls, no external API dependencies\n"
        "• Output ~15 KB submission.csv\n"
        "• Validated against submission_metadata.yaml constraints"
    )
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# ── Slide 9: Technologies Used ──
shape = find_placeholder(prs.slides[8], 'What technologies, frameworks')
if shape:
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "Core Pipeline:\n"
        "• Python 3.10+ — Pure Python scoring (no ML models, no GPU)\n"
        "• JSONL streaming — handles 465 MB without loading into memory entirely\n\n"
        "Backend API:\n"
        "• FastAPI — lightweight REST API for ranking service\n"
        "• Uvicorn — ASGI server\n\n"
        "Frontend:\n"
        "• Vue 3 + Vue Router + Pinia — SPA with PrimeVue UI components\n"
        "• ApexCharts — interactive candidate score visualizations\n"
        "• Vite — fast dev/build tooling\n\n"
        "Dashboard:\n"
        "• Streamlit — quick interactive exploration of results\n"
        "• Plotly — radar charts, histograms, scatter plots\n\n"
        "Version Control:\n"
        "• Git + Git LFS — for large JSONL dataset (465 MB)"
    )
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# ── Slide 10: Submission Assets ──
shape = find_placeholder(prs.slides[9], 'Github')
if shape:
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "GitHub Repository: https://github.com/amalkvr11/india-runs-red-rob-hackathon\n\n"
        "Submission Assets:\n"
        "• submission.csv — top 100 ranked candidates with scores and reasoning\n"
        "• submission_metadata.yaml — scoring dimensions, weights, methodology\n"
        "• This presentation — overview of approach and results\n"
        "• Full source code — ranker pipeline, scorers, honeypot detector, API, frontend"
    )
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# Save
output_path = r'D:\red_rob_hackathon\Idea Submission Template _ Redrob_Filled.pptx'
prs.save(output_path)
print(f'PPT saved to: {output_path}')
