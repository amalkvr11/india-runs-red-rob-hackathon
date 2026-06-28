from pptx import Presentation
from pptx.util import Inches
prs = Presentation(r'D:\red_rob_hackathon\Idea Submission Template _ Redrob.pptx')
slide = prs.slides[6]
for i, shape in enumerate(slide.shapes):
    print(f'Shape {i}: type={shape.shape_type} left={shape.left} top={shape.top} w={shape.width} h={shape.height}')
    if shape.has_text_frame:
        print(f'  text: "{shape.text_frame.text[:100]}"')
